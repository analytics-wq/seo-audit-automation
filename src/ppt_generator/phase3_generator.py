"""Phase 3: PowerPoint Report Generator."""
import logging
import json
from pathlib import Path
from typing import Dict, Any
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class Phase3Generator:
    """Generates Phase 3: Final PowerPoint Output."""

    def __init__(self, phase1_insights: Dict[str, Any],
                 phase2_narrative: Dict[str, Any],
                 template_path: Path = None):
        """Initialize PPT generator.

        Args:
            phase1_insights: Phase 1 analysis results
            phase2_narrative: Phase 2 narrative results
            template_path: Path to PowerPoint template (optional)
        """
        self.phase1 = phase1_insights
        self.phase2 = phase2_narrative
        self.template_path = template_path

        # Merge all data
        self.data = {**phase1_insights, **phase2_narrative}

    def execute(self, output_path: Path) -> Path:
        """Execute Phase 3 PPT generation.

        Args:
            output_path: Path where PPT should be saved

        Returns:
            Path to generated PowerPoint file
        """
        logger.info("=== Starting Phase 3: PowerPoint Report Generation ===")

        # For now, generate a JSON output with all content
        # In production, this would use python-pptx to populate the template
        json_output_path = output_path.parent / f"{output_path.stem}_content.json"

        logger.info(f"Generating structured content output to {json_output_path}")
        self._generate_json_output(json_output_path)

        # Also generate a basic PPT
        ppt_output_path = self._generate_basic_ppt(output_path)

        logger.info(f"=== Phase 3 Complete ===")
        logger.info(f"Content JSON: {json_output_path}")
        logger.info(f"PowerPoint: {ppt_output_path}")

        return ppt_output_path

    def _generate_json_output(self, output_path: Path):
        """Generate JSON file with all structured content."""
        # Convert Pydantic models to dict
        output_data = {}

        for key, value in self.data.items():
            if hasattr(value, 'model_dump'):
                output_data[key] = value.model_dump()
            elif isinstance(value, dict):
                output_data[key] = value
            else:
                output_data[key] = str(value)

        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)

        logger.info(f"Generated JSON content file: {output_path}")

    def _generate_basic_ppt(self, output_path: Path) -> Path:
        """Generate a basic PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(prs)

        # Executive Summary
        self._add_executive_summary_slide(prs)

        # Organic Traffic
        self._add_data_slide(prs, "Organic Traffic Analysis",
                            self.phase1.get('organic_traffic'))

        # Competitive
        self._add_data_slide(prs, "Competitive Benchmarking",
                            self.phase1.get('competitive'))

        # Engagement
        self._add_data_slide(prs, "User Engagement",
                            self.phase1.get('engagement'))

        # Site Health
        self._add_data_slide(prs, "Site Health",
                            self.phase1.get('site_health'))

        # Section Summary
        self._add_summary_slide(prs, "Where You Stand - Summary",
                               self.phase1.get('section_summary_organic'))

        # Meta Tags
        self._add_data_slide(prs, "Meta Tags & On-Page SEO",
                            self.phase1.get('meta_tags'))

        # Keyword Gap
        self._add_data_slide(prs, "Keyword Gap Analysis",
                            self.phase1.get('keyword_gap'))

        # Keyword Intent
        self._add_data_slide(prs, "Keyword Intent Distribution",
                            self.phase1.get('keyword_intent'))

        # Content Summary
        self._add_summary_slide(prs, "Content Gaps - Summary",
                               self.phase1.get('section_summary_content'))

        # Technical SEO
        self._add_data_slide(prs, "Technical SEO Issues",
                            self.phase1.get('technical_seo'))

        # Technical Summary
        self._add_summary_slide(prs, "Technical Gaps - Summary",
                               self.phase1.get('section_summary_technical'))

        # Domain Authority
        self._add_data_slide(prs, "Domain Authority",
                            self.phase1.get('domain_authority'))

        # Authority Summary
        self._add_summary_slide(prs, "Domain Authority - Summary",
                               self.phase1.get('section_summary_authority'))

        # Findings Summary
        self._add_findings_summary_slide(prs)

        # KPIs
        self._add_kpi_slide(prs)

        # Save presentation
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return output_path

    def _add_title_slide(self, prs):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        metadata = self.phase1.get('metadata', {})
        brand_name = metadata.get('brand_name', 'Brand') if hasattr(metadata, 'get') else metadata.brand_name

        title.text = f"SEO Audit\n{brand_name}"
        subtitle.text = datetime.now().strftime("%B %Y")

    def _add_executive_summary_slide(self, prs):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Executive Summary: Path to Digital Visibility"

        # Add text box with summary
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        exec_summary = self.phase2.get('exec_summary', {})

        if hasattr(exec_summary, 'model_dump'):
            summary_dict = exec_summary.model_dump()
        else:
            summary_dict = exec_summary

        for pillar, content in summary_dict.items():
            p = text_frame.add_paragraph()
            p.text = f"{pillar.replace('_', ' ').title()}: {content}"
            p.space_after = Pt(12)

    def _add_data_slide(self, prs, title_text: str, data):
        """Add a data analysis slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = title_text

        if data is None:
            return

        # Add key message
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        if hasattr(data, 'key_message'):
            key_msg = data.key_message
        elif isinstance(data, dict):
            key_msg = data.get('key_message', '')
        else:
            key_msg = ''

        p = text_frame.paragraphs[0]
        p.text = f"Key Finding: {key_msg}"
        p.font.bold = True
        p.font.size = Pt(14)

        # Add observation
        if hasattr(data, 'observation'):
            obs = data.observation
        elif isinstance(data, dict):
            obs = data.get('observation', '')
        else:
            obs = ''

        if obs:
            top = Inches(3)
            textbox2 = slide.shapes.add_textbox(left, top, width, Inches(3))
            text_frame2 = textbox2.text_frame
            p2 = text_frame2.paragraphs[0]
            p2.text = obs
            p2.font.size = Pt(12)

    def _add_summary_slide(self, prs, title_text: str, data):
        """Add a section summary slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = title_text

        if data is None:
            return

        # Create 3-column layout
        left_col = Inches(0.5)
        mid_col = Inches(3.5)
        right_col = Inches(6.5)
        top = Inches(2)
        width = Inches(2.8)
        height = Inches(4)

        # Issues column
        issues_box = slide.shapes.add_textbox(left_col, top, width, height)
        issues_frame = issues_box.text_frame
        p = issues_frame.paragraphs[0]
        p.text = "What is the Issue?"
        p.font.bold = True
        p.font.size = Pt(12)

        issues = data.issues if hasattr(data, 'issues') else data.get('issues', [])
        for issue in issues:
            p = issues_frame.add_paragraph()
            p.text = f"• {issue}"
            p.font.size = Pt(10)

        # Impacts column
        impacts_box = slide.shapes.add_textbox(mid_col, top, width, height)
        impacts_frame = impacts_box.text_frame
        p = impacts_frame.paragraphs[0]
        p.text = "What is the Impact?"
        p.font.bold = True
        p.font.size = Pt(12)

        impacts = data.impacts if hasattr(data, 'impacts') else data.get('impacts', [])
        for impact in impacts:
            p = impacts_frame.add_paragraph()
            p.text = f"• {impact}"
            p.font.size = Pt(10)

        # Actions column
        actions_box = slide.shapes.add_textbox(right_col, top, width, height)
        actions_frame = actions_box.text_frame
        p = actions_frame.paragraphs[0]
        p.text = "Our Next Action"
        p.font.bold = True
        p.font.size = Pt(12)

        actions = data.actions if hasattr(data, 'actions') else data.get('actions', [])
        for action in actions:
            p = actions_frame.add_paragraph()
            p.text = f"• {action}"
            p.font.size = Pt(10)

    def _add_findings_summary_slide(self, prs):
        """Add overall findings summary slide."""
        findings = self.phase2.get('findings_summary')
        if findings:
            self._add_data_slide(prs, "SEO Audit Findings Summary", findings)

    def _add_kpi_slide(self, prs):
        """Add KPI & Benchmark slide."""
        kpi = self.phase1.get('kpi')
        if kpi:
            self._add_data_slide(prs, "KPI & Benchmark", kpi)
