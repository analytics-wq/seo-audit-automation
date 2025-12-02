#!/usr/bin/env python3
"""
Create a standardized PowerPoint template matching brand_template_ppt_report.pdf
This creates a base .pptx file with proper layouts and placeholder text boxes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_base_template():
    """
    Create a PowerPoint template with standardized layouts
    matching the brand_template_ppt_report.pdf structure
    """
    prs = Presentation()

    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating SEO Audit PowerPoint Template...")
    print("This template matches brand_template_ppt_report.pdf structure")

    # Define brand colors (can be customized)
    BRAND_PRIMARY = RGBColor(0, 32, 96)      # Dark blue
    BRAND_SECONDARY = RGBColor(156, 39, 176)  # Purple
    BRAND_ACCENT = RGBColor(255, 193, 7)      # Yellow/Gold
    TEXT_GRAY = RGBColor(100, 100, 100)

    # Store colors in presentation for reference
    template_data = {
        "brand_colors": {
            "primary": BRAND_PRIMARY,
            "secondary": BRAND_SECONDARY,
            "accent": BRAND_ACCENT,
            "text_gray": TEXT_GRAY
        },
        "fonts": {
            "heading": "Calibri",
            "body": "Calibri"
        }
    }

    # SLIDE 0: Cover Page
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Logo placeholder (top right)
    logo_box = slide.shapes.add_textbox(Inches(8), Inches(0.5), Inches(1.5), Inches(0.8))
    logo_frame = logo_box.text_frame
    logo_p = logo_frame.paragraphs[0]
    logo_p.text = "{Logo}"
    logo_p.font.size = Pt(10)
    logo_p.font.color.rgb = TEXT_GRAY
    logo_p.alignment = PP_ALIGN.CENTER

    # Brand name placeholder (center)
    brand_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    brand_frame = brand_box.text_frame
    brand_frame.word_wrap = True
    brand_p = brand_frame.paragraphs[0]
    brand_p.text = "SEO Audit Proposal\n{Brand}"
    brand_p.font.size = Pt(44)
    brand_p.font.bold = True
    brand_p.font.color.rgb = BRAND_PRIMARY
    brand_p.alignment = PP_ALIGN.CENTER

    # Date placeholder
    date_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(0.5))
    date_frame = date_box.text_frame
    date_p = date_frame.paragraphs[0]
    date_p.text = "{Current month}"
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = TEXT_GRAY
    date_p.alignment = PP_ALIGN.CENTER

    print("  ✓ Slide 0: Cover Page created")

    # SLIDES 1-3: Setup slides (using standard Title + Content layout)
    for i, title in enumerate(["SEO Audit Overview", "Planning Phase", "Index"], 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_PRIMARY
        print(f"  ✓ Slide {i}: {title} created")

    # SLIDE 4: Executive Summary (4-column layout)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Path to Digital Visibility"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_PRIMARY

    # Add 4 placeholder columns
    column_width = Inches(2)
    column_start = Inches(0.5)
    column_y = Inches(2)

    for i, col_title in enumerate(["General Overview", "Content SEO", "Technical SEO", "Domain Authority"]):
        x_pos = column_start + (i * Inches(2.4))

        # Column title
        title_box = slide.shapes.add_textbox(x_pos, column_y, column_width, Inches(0.4))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = col_title
        title_p.font.size = Pt(12)
        title_p.font.bold = True
        title_p.font.color.rgb = BRAND_SECONDARY

        # Column content placeholder
        content_box = slide.shapes.add_textbox(x_pos, column_y + Inches(0.5), column_width, Inches(3.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_p = content_frame.paragraphs[0]
        content_p.text = f"{{Summary {col_title.lower()}}}"
        content_p.font.size = Pt(10)
        content_p.font.color.rgb = TEXT_GRAY

    print("  ✓ Slide 4: Executive Summary created")

    # SLIDES 5-25: Data slides with standard layout
    data_slides = [
        (5, "Where {Brand} Stands Today?", "divider"),
        (6, "Organic Traffic Analysis", "data"),
        (7, "Competitive Benchmarking", "data"),
        (8, "User Engagement", "data"),
        (9, "Site Health", "data"),
        (10, "Summary - Where You Stand", "summary"),
        (11, "Content Visibility Gaps & Insights", "divider"),
        (12, "Meta Tags & Heading", "data"),
        (13, "Keyword Gap Analysis", "data"),
        (14, "Keyword Intent Distribution", "data"),
        (15, "Summary - Content Gaps", "summary"),
        (16, "Technical Gaps Limiting Growth", "divider"),
        (17, "Technical SEO Issues", "data"),
        (18, "Summary - Technical Gaps", "summary"),
        (19, "Understanding Domain Authority", "divider"),
        (20, "Domain Authority Assessment", "data"),
        (21, "Summary - Domain Authority", "summary"),
        (22, "Steps to Improve Your SEO", "divider"),
        (23, "SEO Audit Findings Summary", "data"),
        (24, "KPI Targets & Benchmarks", "data"),
        (25, "Thank You", "closing")
    ]

    for slide_num, title, slide_type in data_slides:
        if slide_type == "divider":
            # Section divider - minimal design
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            title_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_p = title_frame.paragraphs[0]
            title_p.text = title
            title_p.font.size = Pt(36)
            title_p.font.bold = True
            title_p.font.color.rgb = BRAND_PRIMARY
            title_p.alignment = PP_ALIGN.CENTER

        elif slide_type == "summary":
            # Summary slide - 3-box layout
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_PRIMARY

            # Add placeholder for summary content
            summary_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
            summary_frame = summary_box.text_frame
            summary_frame.word_wrap = True
            summary_p = summary_frame.paragraphs[0]
            summary_p.text = "{Summary content - Issue | Impact | Action}"
            summary_p.font.size = Pt(12)
            summary_p.font.color.rgb = TEXT_GRAY

        elif slide_type == "closing":
            # Thank you slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
            thank_you_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
            thank_you_frame = thank_you_box.text_frame
            thank_you_frame.word_wrap = True
            thank_you_p = thank_you_frame.paragraphs[0]
            thank_you_p.text = "Thank You"
            thank_you_p.font.size = Pt(44)
            thank_you_p.font.bold = True
            thank_you_p.font.color.rgb = BRAND_PRIMARY
            thank_you_p.alignment = PP_ALIGN.CENTER

        else:  # data slide
            # Standard data slide with key highlight box
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_PRIMARY

            # Priority badge placeholder (top right)
            badge_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.2), Inches(1), Inches(0.5))
            badge_frame = badge_box.text_frame
            badge_p = badge_frame.paragraphs[0]
            badge_p.text = "{Priority}"
            badge_p.font.size = Pt(18)
            badge_p.font.bold = True
            badge_p.alignment = PP_ALIGN.CENTER

            # Key Highlight section
            highlight_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            highlight_frame = highlight_box.text_frame
            highlight_frame.word_wrap = True

            # "KEY INSIGHT" label
            label_p = highlight_frame.paragraphs[0]
            label_p.text = "KEY INSIGHT"
            label_p.font.size = Pt(10)
            label_p.font.bold = True
            label_p.font.color.rgb = TEXT_GRAY

            # Key message placeholder
            key_msg_p = highlight_frame.add_paragraph()
            key_msg_p.text = "{Key Highlight}"
            key_msg_p.font.size = Pt(16)
            key_msg_p.font.bold = True
            key_msg_p.font.color.rgb = RGBColor(0, 0, 0)
            key_msg_p.space_before = Pt(10)

            # Observation section
            obs_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(3.5))
            obs_frame = obs_box.text_frame
            obs_frame.word_wrap = True
            obs_p = obs_frame.paragraphs[0]
            obs_p.text = "{Observation & Analysis}\n\n{Data visualizations and detailed analysis go here}"
            obs_p.font.size = Pt(12)
            obs_p.font.color.rgb = TEXT_GRAY

        print(f"  ✓ Slide {slide_num}: {title} created")

    # Save the template
    template_path = "/home/user/seo-audit-automation/template_base.pptx"
    prs.save(template_path)

    print(f"\n✅ Template created: {template_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   File size: {Path(template_path).stat().st_size / 1024:.1f} KB")
    print("\nNext: Use this template in generate_seo_report.py")

    return template_path

if __name__ == "__main__":
    from pathlib import Path
    create_base_template()
